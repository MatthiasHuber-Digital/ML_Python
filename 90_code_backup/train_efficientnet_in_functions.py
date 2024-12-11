import os
import torch
from torch.utils.tensorboard import SummaryWriter
from torch.utils.data import DataLoader
from torch import optim
import torch.nn as nn
from torchvision import datasets
from sklearn.preprocessing import LabelBinarizer
from sklearn.metrics import f1_score, precision_score, recall_score, average_precision_score
from efficientnet_pytorch import EfficientNet
from efficientnet_pytorch import EfficientNet
from efficient_net.options import OptionsModelType
from efficient_net.utils import load_efficient_net_model
from efficient_net.split_dataset_kfoldcrossval import get_kfold_directory_list
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
import numpy as np
import os
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from sklearn.preprocessing import LabelBinarizer
from sklearn.metrics import (
    f1_score,
    precision_score,
    precision_recall_curve,
    recall_score,
    average_precision_score,
    roc_auc_score,
    roc_curve,
    auc,
)
import torch
from torch import nn
import torch.optim as optim
from torch.utils.data import DataLoader
from torch.utils.data.dataloader import default_collate
from torch.utils.tensorboard import SummaryWriter
import torchvision.transforms.functional as F
from torchvision import datasets, transforms
from tqdm import tqdm


# Path and settings
PATH_LOAD_FROM_DISK = "/home/matthias/workspace/Coding/00_vista_medizina/10_weights/efficient_net/2024-12-02_bf_kaggle/efficientnet-b4_BEST.pth"
# "/home/matthias/workspace/Coding/00_vista_medizina/10_weights/efficient_net/2024-12-02_bf_kaggle/efficientnet-b5_epoch30_BEST.pth"
# "/home/matthias/workspace/Coding/00_vista_medizina/10_weights/efficient_net/2024-12-02_bf_kaggle/efficientnet-b4_BEST.pth"
PATH_DATASET = "/home/matthias/workspace/Coding/00_vista_medizina/00_data/2024-11-08/kaggle_bone_fracture_detection_small_2024-11-08/bf_orig_bgr_classification_7cl_kfold/w_folds"
PATH_OUTPUT_DIR = "/home/matthias/workspace/Coding/00_vista_medizina/vista_bone_frac/efficient_net/output_7cl_5fold_3runs"

LR_INITIAL = 5e-6  # best: 1e-6
LR_PLATEAU_REDUCER = True  # const lr: False
LR_PLAT_FACTOR = 0.5  # FLOAT, pytorch: 0.1
LR_PLAT_PATIENCE = 5  # INT, pytorch:3
MODEL_TYPE = OptionsModelType.B4  #  #
if MODEL_TYPE == OptionsModelType.B5:
    BATCH_SIZE = 10
elif MODEL_TYPE == OptionsModelType.B4:
    BATCH_SIZE = 22
# K_FOLDS = 5  # Number of folds for cross-validation
NUM_EPOCHS = 31  # 31  # Number of epochs to train
NUM_RUNS = 1  # 3  # Number of runs for random splits
LIST_SAVE_EPOCHS = [10, 20, 30]
SAVE_CURRENT_BEST_MODEL = True  # False

# Device setup
device = torch.device("cuda" if torch.cuda.is_available() else "cpu")


# Transforms for training and validation
train_transforms = transforms.Compose(
    [
        ResizeWithAspectRatioAndPadding(MODEL_TYPE.get_resolution()),
        transforms.RandomHorizontalFlip(),
        transforms.RandomVerticalFlip(),
        transforms.RandomRotation(20),
        transforms.ToTensor(),
        transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
    ]
)

val_transforms = transforms.Compose(
    [
        ResizeWithAspectRatioAndPadding(MODEL_TYPE.get_resolution()),
        transforms.ToTensor(),
        transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
    ]
)


def plot_roc_curve(
    y_true,
    y_pred_probs,
    save_object,
    title="ROC Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # Convert labels to one-hot format for multiclass classification
    lb = LabelBinarizer()
    y_true_bin = lb.fit_transform(y_true)

    if not isinstance(y_true_bin, np.ndarray):
        y_true_bin = np.array(y_true_bin)
    if not isinstance(y_pred_probs, np.ndarray):
        y_pred_probs = np.array(y_pred_probs)

    # Ensure y_pred_probs is a NumPy array
    y_pred_probs = np.array(y_pred_probs)

    # Check if y_pred_probs is 2D (for multiclass)
    if y_pred_probs.ndim == 1:
        # If it's 1D, expand it to 2D for binary classification (one class vs rest)
        y_pred_probs = np.expand_dims(y_pred_probs, axis=1)

    # Plot ROC curve for each class
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    n_classes = y_true_bin.shape[1]
    fpr, tpr, roc_auc = {}, {}, {}

    for i in range(n_classes):
        fpr[i], tpr[i], _ = roc_curve(y_true_bin[:, i], y_pred_probs[:, i])
        roc_auc[i] = auc(fpr[i], tpr[i])
        plt.plot(fpr[i], tpr[i], label=f"Class {lb.classes_[i]} (AUC = {roc_auc[i]:.2f})")

    plt.plot([0, 1], [0, 1], color="gray", linestyle="--")  # Diagonal line (chance level)
    plt.xlabel("False Positive Rate")
    plt.ylabel("True Positive Rate")
    plt.title(title)
    plt.legend(loc="lower right")

    # Save the plot to the provided PDF
    save_object.savefig()
    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_precision_recall_curve(
    y_true,
    y_pred_probs,
    save_object,
    title="Precision-Recall Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # Convert labels to one-hot format for multiclass classification
    lb = LabelBinarizer()
    y_true_bin = lb.fit_transform(y_true)

    if not isinstance(y_true_bin, np.ndarray):
        y_true_bin = np.array(y_true_bin)
    if not isinstance(y_pred_probs, np.ndarray):
        y_pred_probs = np.array(y_pred_probs)

    # Ensure y_pred_probs is a NumPy array
    y_pred_probs = np.array(y_pred_probs)

    # Check if y_pred_probs is 2D (for multiclass)
    if y_pred_probs.ndim == 1:
        # If it's 1D, expand it to 2D for binary classification (one class vs rest)
        y_pred_probs = np.expand_dims(y_pred_probs, axis=1)

    # Plot Precision-Recall curve for each class
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    n_classes = y_true_bin.shape[1]

    for i in range(n_classes):
        precision, recall, _ = precision_recall_curve(y_true_bin[:, i], y_pred_probs[:, i])
        avg_precision = average_precision_score(y_true_bin[:, i], y_pred_probs[:, i])
        plt.plot(recall, precision, label=f"Class {lb.classes_[i]} (AP = {avg_precision:.2f})")

    plt.xlabel("Recall")
    plt.ylabel("Precision")
    plt.title(title)
    plt.legend(loc="lower left")

    # Save the plot to the provided PDF
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)

    plt.close()


def plot_loss_curve(
    losses,
    save_object,
    title="Loss Curve",  # tensorboard_writer: SummaryWriter = None
):
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    plt.plot(losses, color="red", lw=2)
    plt.xlabel("Epochs")
    plt.ylabel("Loss")
    plt.title(title)
    save_object.savefig()
    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_class_distribution(
    list_class_occurences,
    save_object,
    title="Class Distribution",
    # tensorboard_writer: SummaryWriter = None,
):
    # sns.barplot(x=list(class_counts.keys()), y=list(class_counts.values()))
    # fixed bin size
    bins = np.arange(0, max(list_class_occurences), 1)  # fixed bin size

    plt.xlim([min(list_class_occurences), max(list_class_occurences)])
    # fig = plt.hist(list_class_occurences, bins=bins, alpha=0.5)
    plt.hist(list_class_occurences, bins=bins, alpha=0.5)
    plt.xlabel("Classes")
    plt.ylabel("Count")
    plt.title(title)
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_loss_accuracy_curve(
    list_train_losses,
    list_val_losses,
    list_train_accuracies,
    list_val_accuracies,
    save_object,
    title="Loss & Accuracy Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # fig = plt.figure(figsize=(10, 5))
    plt.figure(figsize=(10, 5))
    plt.title(title)
    plt.subplot(1, 2, 1)
    plt.plot(range(NUM_EPOCHS), list_train_losses, label="Training Loss")
    plt.plot(range(NUM_EPOCHS), list_val_losses, label="Validation Loss")
    plt.xlabel("Epochs")
    plt.ylabel("Loss")
    plt.legend()

    plt.subplot(1, 2, 2)
    plt.plot(range(NUM_EPOCHS), list_train_accuracies, label="Training Accuracy")
    plt.plot(range(NUM_EPOCHS), list_val_accuracies, label="Validation Accuracy")
    plt.xlabel("Epochs")
    plt.ylabel("Accuracy (%)")
    plt.legend()

    plt.tight_layout()
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


# Helper function to create a PowerPoint presentation slide
def add_ppt_slide(ppt, title, content, image_path=None):
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Use the 'Title and Content' layout

    # Set slide title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Add content (text or image)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(9))  # Add image
    else:
        content_box = slide.shapes.placeholders[1]  # Text box
        content_box.text = content


# Collate function
def custom_collate_fn(batch):
    # Ensure all images in the batch have the same size
    return default_collate(batch)


# Model transforms
class ResizeWithAspectRatioAndPadding:
    def __init__(self, target_size):
        self.target_size = target_size

    def __call__(self, img):
        # Get the original dimensions of the image
        width, height = img.size
        aspect_ratio = width / height

        # Compute the new dimensions while maintaining the aspect ratio
        if aspect_ratio > 1:
            new_width = self.target_size
            new_height = int(self.target_size / aspect_ratio)
        else:
            new_height = self.target_size
            new_width = int(self.target_size * aspect_ratio)

        # Resize the image
        img = F.resize(img, (new_height, new_width))

        # Now, pad the image to match the target size
        padding = (0, 0, self.target_size - new_width, self.target_size - new_height)
        img = F.pad(img, padding, fill=0)
        return img


dict_kfold_dirs = get_kfold_directory_list(path_image_dir=PATH_DATASET)
list_paths_fold_train_dirs = dict_kfold_dirs["list_paths_fold_train_dirs"]
list_paths_fold_val_dirs = dict_kfold_dirs["list_paths_fold_val_dirs"]


def create_run_directory(run, output_dir):
    run_output_dir = os.path.join(output_dir, f"run_{run}")
    os.makedirs(run_output_dir, exist_ok=True)
    return run_output_dir


def create_fold_directory(run_output_dir, fold_idx):
    fold_output_dir = os.path.join(run_output_dir, f"fold_{fold_idx}")
    os.makedirs(fold_output_dir, exist_ok=True)
    return fold_output_dir


def create_dataloaders(train_dir, val_dir, batch_size, transform_train, transform_val, collate_fn):
    train_dataset = datasets.ImageFolder(train_dir, transform=transform_train)
    val_dataset = datasets.ImageFolder(val_dir, transform=transform_val)

    train_loader = DataLoader(
        train_dataset, batch_size=batch_size, shuffle=True, num_workers=4, collate_fn=collate_fn
    )
    val_loader = DataLoader(
        val_dataset, batch_size=batch_size, shuffle=False, num_workers=4, collate_fn=collate_fn
    )

    return train_loader, val_loader, train_dataset, val_dataset


def load_model(path_load_from_disk, model_type, num_classes, device):
    if path_load_from_disk:
        model = load_efficient_net_model(path_load_from_disk, model_type, num_classes)
    else:
        model = EfficientNet.from_pretrained(model_type, num_classes=num_classes)
    model.train()
    model.to(device)
    return model


def initialize_optimizer_and_scheduler(
    model, lr_initial, lr_plateau_reducer, lr_plat_factor, lr_plat_patience
):
    criterion = nn.CrossEntropyLoss()
    optimizer = optim.Adam(model.parameters(), lr=lr_initial)

    scaler = torch.amp.GradScaler()

    if lr_plateau_reducer:
        scheduler = torch.optim.lr_scheduler.ReduceLROnPlateau(
            optimizer, mode="min", factor=lr_plat_factor, patience=lr_plat_patience, verbose=True
        )
    else:
        scheduler = None

    return criterion, optimizer, scaler, scheduler


def train_one_epoch(model, train_loader, criterion, optimizer, scaler, device):
    model.train()
    running_loss = 0.0
    correct_preds = 0
    total_preds = 0
    all_preds = []
    all_labels = []
    all_pred_probs = []

    for images, labels in train_loader:
        images, labels = images.to(device), labels.to(device)
        optimizer.zero_grad()

        with torch.amp.autocast(device_type=str(device)):  # Mixed precision training
            output = model(images)
            loss = criterion(output, labels)

        scaler.scale(loss).backward()
        scaler.step(optimizer)
        scaler.update()

        running_loss += loss.item()
        _, predicted = torch.max(output, 1)
        correct_preds += (predicted == labels).sum().item()
        total_preds += labels.size(0)

        all_preds.extend(predicted.cpu().numpy())
        all_labels.extend(labels.cpu().numpy())

        # Collect predicted probabilities for average_precision_score
        pred_probs = torch.softmax(output, dim=1)
        all_pred_probs.extend(pred_probs.cpu().detach().numpy())

    train_loss = running_loss / len(train_loader)
    train_acc = 100 * correct_preds / total_preds

    # Calculate metrics
    train_f1 = f1_score(all_labels, all_preds, average="weighted")
    train_precision = precision_score(all_labels, all_preds, average="weighted")
    train_recall = recall_score(all_labels, all_preds, average="weighted")
    train_average_precision = average_precision_score(
        LabelBinarizer().fit_transform(all_labels), all_pred_probs, average="weighted"
    )

    return train_loss, train_acc, train_f1, train_precision, train_recall, train_average_precision


def validate_one_epoch(model, val_loader, criterion, device):
    model.eval()
    running_loss = 0.0
    correct_preds = 0
    total_preds = 0
    all_preds = []
    all_labels = []
    all_pred_probs = []

    with torch.no_grad():
        for images, labels in val_loader:
            images, labels = images.to(device), labels.to(device)
            output = model(images)
            loss = criterion(output, labels)
            running_loss += loss.item()
            _, predicted = torch.max(output, 1)
            correct_preds += (predicted == labels).sum().item()
            total_preds += labels.size(0)

            all_preds.extend(predicted.cpu().numpy())
            all_labels.extend(labels.cpu().numpy())

            # Collect predicted probabilities for average_precision_score
            pred_probs = torch.softmax(output, dim=1)
            all_pred_probs.extend(pred_probs.cpu().detach().numpy())

    val_loss = running_loss / len(val_loader)
    val_acc = 100 * correct_preds / total_preds

    # Calculate metrics
    val_f1 = f1_score(all_labels, all_preds, average="weighted")
    val_precision = precision_score(all_labels, all_preds, average="weighted")
    val_recall = recall_score(all_labels, all_preds, average="weighted")
    val_average_precision = average_precision_score(
        LabelBinarizer().fit_transform(all_labels), all_pred_probs, average="weighted"
    )

    return val_loss, val_acc, val_f1, val_precision, val_recall, val_average_precision


def log_to_tensorboard(writer, epoch, train_metrics, val_metrics):
    train_loss, train_acc, train_f1, train_precision, train_recall, train_avg_precision = (
        train_metrics
    )
    val_loss, val_acc, val_f1, val_precision, val_recall, val_avg_precision = val_metrics

    writer.add_scalar("Loss/train", train_loss, epoch)
    writer.add_scalar("Accuracy/train", train_acc, epoch)
    writer.add_scalar("F1/train", train_f1, epoch)
    writer.add_scalar("Precision/train", train_precision, epoch)
    writer.add_scalar("Recall/train", train_recall, epoch)
    writer.add_scalar("Average Precision/train", train_avg_precision, epoch)

    writer.add_scalar("Loss/val", val_loss, epoch)
    writer.add_scalar("Accuracy/val", val_acc, epoch)
    writer.add_scalar("F1/val", val_f1, epoch)
    writer.add_scalar("Precision/val", val_precision, epoch)
    writer.add_scalar("Recall/val", val_recall, epoch)
    writer.add_scalar("Average Precision/val", val_avg_precision, epoch)


def save_best_model(model, fold_output_dir, epoch, val_acc, highest_val_acc):
    if val_acc > highest_val_acc:
        highest_val_acc = val_acc
        print("==> Model saved with highest validation accuracy!")
        torch.save(
            model.state_dict(), os.path.join(fold_output_dir, f"model_epoch_{epoch}_BEST.pth")
        )
    return highest_val_acc


def save_model_at_specified_epochs(model, fold_output_dir, epoch, list_save_epochs):
    if epoch in list_save_epochs:
        print("==> Saving model to disk...")
        torch.save(
            model.state_dict(), os.path.join(fold_output_dir, f"model_epoch_{epoch}_SCHED.pth")
        )


def main_training_loop(
    num_runs, list_paths_fold_train_dirs, list_paths_fold_val_dirs, output_dir, device
):
    for run in range(num_runs):
        print(f"\nTRAIN RUN NUMBER: {run}/{num_runs-1}")
        run_output_dir = create_run_directory(run, output_dir)
        writer = SummaryWriter(os.path.join(run_output_dir, "tensorboard"))

        with PdfPages(os.path.join(run_output_dir, f"training_report_run_{run}.pdf")) as pdf_pages:
            for idx_fold, (train_dir, val_dir) in enumerate(
                zip(list_paths_fold_train_dirs, list_paths_fold_val_dirs)
            ):
                print(f"\n----Fold {idx_fold}----------------------------")
                fold_output_dir = create_fold_directory(run_output_dir, idx_fold)
                train_loader, val_loader, train_dataset, val_dataset = create_dataloaders(
                    train_dir,
                    val_dir,
                    BATCH_SIZE,
                    train_transforms,
                    val_transforms,
                    custom_collate_fn,
                )

                model = load_model(
                    PATH_LOAD_FROM_DISK, MODEL_TYPE, len(train_dataset.classes), device
                )
                criterion, optimizer, scaler, scheduler = initialize_optimizer_and_scheduler(
                    model, LR_INITIAL, LR_PLATEAU_REDUCER, LR_PLAT_FACTOR, LR_PLAT_PATIENCE
                )

                highest_val_acc = 0.0
                for epoch in range(NUM_EPOCHS):
                    print(f"Epoch {epoch}/{NUM_EPOCHS-1}")

                    # Train and validate for the epoch
                    train_metrics = train_one_epoch(
                        model, train_loader, criterion, optimizer, scaler, device
                    )
                    val_metrics = validate_one_epoch(model, val_loader, criterion, device)

                    # Log metrics
                    log_to_tensorboard(writer, epoch, train_metrics, val_metrics)

                    # Save best model and models at specified epochs
                    highest_val_acc = save_best_model(
                        model, fold_output_dir, epoch, val_metrics[1], highest_val_acc
                    )
                    save_model_at_specified_epochs(model, fold_output_dir, epoch, LIST_SAVE_EPOCHS)

                    print(f"Train Acc: {train_metrics[1]}%, Val Acc: {val_metrics[1]}%")

                # Optionally save plots to PDF here

        writer.close()
        print(f"Finished run {run + 1}/{num_runs}.")


if __name__ == "__main__":
    main_training_loop(
        NUM_RUNS: int=1, 
        list_paths_fold_train_dirs: str, 
        list_paths_fold_val_dirs: str, 
        PATH_OUTPUT_DIR: str, 
        device: str,
    )
