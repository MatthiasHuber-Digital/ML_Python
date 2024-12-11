import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from sklearn.metrics import roc_curve, auc, average_precision_score
import numpy as np
import seaborn as sns
from sklearn.metrics import precision_recall_curve
from sklearn.preprocessing import LabelBinarizer
from pptx import Presentation
from pptx.util import Inches
import os

# Function to generate and save a plot for a given metric (e.g., ROC, precision-recall, loss)
def plot_roc_curve(y_true, y_pred, pdf_pages, title="ROC Curve"):
    fpr, tpr, _ = roc_curve(y_true, y_pred)
    roc_auc = auc(fpr, tpr)
    plt.figure(figsize=(8, 6))
    plt.plot(fpr, tpr, color='blue', lw=2, label=f'ROC curve (AUC = {ro32c_auc:.2f})')
    plt.plot([0, 1], [0, 1], color='gray', linestyle='--')
    plt.xlabel('False Positive Rate')
    plt.ylabel('True Positive Rate')
    plt.title(title)
    plt.legend(loc="lower right")
    pdf_pages.savefig()  # Save the figure to the PDF
    plt.close()

def plot_precision_recall_curve(y_true, y_pred, pdf_pages, title="Precision-Recall Curve"):
    precision, recall, _ = precision_recall_curve(y_true, y_pred)
    plt.figure(figsize=(8, 6))
    plt.plot(recall, precision, color='blue', lw=2)
    plt.xlabel('Recall')
    plt.ylabel('Precision')
    plt.title(title)
    pdf_pages.savefig()
    plt.close()

def plot_loss_curve(losses, pdf_pages, title="Loss Curve"):
    plt.figure(figsize=(8, 6))
    plt.plot(losses, color='red', lw=2)
    plt.xlabel('Epochs')
    plt.ylabel('Loss')
    plt.title(title)
    pdf_pages.savefig()
    plt.close()

def plot_class_distribution(class_counts, pdf_pages, title="Class Distribution"):
    sns.barplot(x=list(class_counts.keys()), y=list(class_counts.values()))
    plt.xlabel('Classes')
    plt.ylabel('Count')
    plt.title(title)
    pdf_pages.savefig()
    plt.close()


# Function to generate a PDF report for training results
def generate_pdf_report(run_results):
    with PdfPages('training_report.pdf') as pdf_pages:
        for run in run_results:
            # Save metrics and scalar explanations
            for fold, results in run['fold_results'].items():
                # Generate fold-based class distribution statistic plot
                plot_class_distribution(results['class_distribution'], pdf_pages, title=f"Class Distribution for Fold {fold}")
                
                # Generate and save loss curve plot
                plot_loss_curve(results['losses'], pdf_pages, title=f"Loss Curve for Fold {fold}")
                
                # Generate and save ROC curve plot
                plot_roc_curve(results['y_true'], results['y_pred'], pdf_pages, title=f"ROC Curve for Fold {fold}")
                
                # Generate and save Precision-Recall curve
                plot_precision_recall_curve(results['y_true'], results['y_pred'], pdf_pages, title=f"Precision-Recall Curve for Fold {fold}")
                
                # Additional plots for Average Precision (AP) and mAP can be added in a similar manner




import os
import torch
import torch.optim as optim
import matplotlib.pyplot as plt
from torch.utils.data import DataLoader
from torch import nn
from torchvision import datasets, transforms
from tqdm import tqdm
from sklearn.metrics import (
    f1_score,
    precision_score,
    recall_score,
    average_precision_score,
    roc_curve,
    auc
)
from efficientnet_pytorch import EfficientNet
from efficient_net.utils import load_efficient_net_model
from efficient_net.split_dataset_kfoldcrossval import get_kfold_directory_list
from torch.utils.tensorboard import SummaryWriter
from matplotlib.backends.backend_pdf import PdfPages
import numpy as np
import seaborn as sns

# Paths and settings
PATH_LOAD_FROM_DISK = "/path/to/your/model.pth"
PATH_DATASET = "/path/to/your/dataset"
PATH_OUTPUT_DIR = "/path/to/save/output"
LR_INITIAL = 5e-6
MODEL_TYPE = OptionsModelType.B5
BATCH_SIZE = 10
K_FOLDS = 5
NUM_EPOCHS = 31
NUM_RUNS = 3

device = torch.device("cuda" if torch.cuda.is_available() else "cpu")



# Helper function to create a PowerPoint presentation slide
def add_ppt_slide(ppt, title, content, image_path=None):
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Use the 'Title and Content' layout

    # Set slide title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Add content (text or image)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(9))  # Add image
    else:
        content_box = slide.shapes.placeholders[1]  # Text box
        content_box.text = content


val_all_pred_probs = []  # To store predicted probabilities for average_precision_score
# Training and evaluation loop

for run in range(NUM_RUNS):
    print(f"\nTRAIN RUN NUMBER: {run}/{NUM_RUNS - 1}")

    run_output_dir = os.path.join(PATH_OUTPUT_DIR, f"run_{run}")
    os.makedirs(run_output_dir, exist_ok=True)
    writer = SummaryWriter(os.path.join(run_output_dir, "tensorboard"))

    # Create a PowerPoint presentation object
    ppt = Presentation()


    with PdfPages(os.path.join(run_output_dir, f"training_report_run_{run}.pdf")) as pdf_pages:
        # Loop over each fold
        for dict_path_dir in list_paths_kfold_dirs:
            idx_fold = dict_path_dir["idx_fold"]
            path_kfold_train_dir = dict_path_dir["train"]
            path_kfold_val_dir = dict_path_dir["val"]
            print(f"\n----Fold {idx_fold}----------------------------")

            fold_output_dir = os.path.join(run_output_dir, f"fold_{idx_fold}")
            os.makedirs(fold_output_dir, exist_ok=True)

            # Load datasets and create dataloaders
            kfold_train_dataset = datasets.ImageFolder(path_kfold_train_dir, transform=train_transforms)
            kfold_val_dataset = datasets.ImageFolder(path_kfold_val_dir, transform=val_transforms)

            train_fold_loader = DataLoader(kfold_train_dataset, batch_size=BATCH_SIZE, shuffle=True)
            val_fold_loader = DataLoader(kfold_val_dataset, batch_size=BATCH_SIZE, shuffle=False)

            model = load_efficient_net_model(PATH_LOAD_FROM_DISK, MODEL_TYPE, num_classes=len(kfold_train_dataset.classes))
            model.to(device)

            criterion = nn.CrossEntropyLoss()
            optimizer = optim.Adam(model.parameters(), lr=LR_INITIAL)

            train_losses, val_losses = [], []
            train_accuracies, val_accuracies = [], []

            highest_val_acc = 0.0

            # Training loop
            for epoch in range(NUM_EPOCHS):
                model.train()
                running_loss = 0.0
                correct_preds = 0
                total_preds = 0
                all_preds = []
                all_labels = []

                for images, labels in tqdm(train_fold_loader):
                    images, labels = images.to(device), labels.to(device)
                    optimizer.zero_grad()

                    output = model(images)
                    loss = criterion(output, labels)

                    loss.backward()
                    optimizer.step()

                    running_loss += loss.item()
                    _, predicted = torch.max(output, 1)
                    correct_preds += (predicted == labels).sum().item()
                    total_preds += labels.size(0)

                    all_preds.extend(predicted.cpu().numpy())
                    all_labels.extend(labels.cpu().numpy())

                    # Collect predicted probabilities for average_precision_score
                    pred_probs = torch.softmax(output, dim=1)  # Apply softmax to get probabilities for each class
                    all_pred_probs.extend(pred_probs.cpu().numpy())


                # Calculate F1, Precision, Recall, and Average Precision using the correct inputs
                f1 = f1_score(all_labels, all_preds, average="weighted")
                precision = precision_score(all_labels, all_preds, average="weighted")
                recall = recall_score(all_labels, all_preds, average="weighted")
                
                # Average Precision: Use the probabilities, not just predicted classes
                average_precision = average_precision_score(all_labels, all_pred_probs, average="weighted", multi_class="ovr")

                # Log metrics to TensorBoard
                writer.add_scalar("Loss/train", train_loss, epoch)
                writer.add_scalar("Accuracy/train", train_acc, epoch)
                writer.add_scalar("F1/train", f1, epoch)
                writer.add_scalar("Precision/train", precision, epoch)
                writer.add_scalar("Recall/train", recall, epoch)
                writer.add_scalar("Average Precision/train", average_precision, epoch)

                # Validation loop
                model.eval()
                val_loss, val_correct_preds, val_total_preds = 0.0, 0, 0
                val_all_preds, val_all_labels = [], []
                val_all_pred_probs = []  # To store predicted probabilities for average_precision_score


                with torch.no_grad():
                    for images, labels in tqdm(val_fold_loader):
                        images, labels = images.to(device), labels.to(device)
                        outputs = model(images)
                        loss = criterion(outputs, labels)
                        val_loss += loss.item()
                        _, predicted = torch.max(outputs, 1)
                        val_correct_preds += (predicted == labels).sum().item()
                        val_total_preds += labels.size(0)

                        val_all_preds.extend(predicted.cpu().numpy())
                        val_all_labels.extend(labels.cpu().numpy())
                        
                        # Collect predicted probabilities for average_precision_score
                        pred_probs = torch.softmax(outputs, dim=1)  # Apply softmax to get probabilities for each class
                        val_all_pred_probs.extend(pred_probs.cpu().numpy())

                # Calculate validation metrics
                val_loss = val_loss / len(val_fold_loader)
                val_acc = 100 * val_correct_preds / val_total_preds
                val_f1 = f1_score(val_all_labels, val_all_preds, average="weighted")
                val_precision = precision_score(val_all_labels, val_all_preds, average="weighted")
                val_recall = recall_score(val_all_labels, val_all_preds, average="weighted")
                
                # Average Precision (using probabilities)
                val_avg_precision = average_precision_score(val_all_labels, val_all_pred_probs, average="weighted", multi_class="ovr")


                # Log validation metrics to TensorBoard
                writer.add_scalar("Loss/val", val_loss, epoch)
                writer.add_scalar("Accuracy/val", val_acc, epoch)
                writer.add_scalar("F1/val", val_f1, epoch)
                writer.add_scalar("Precision/val", val_precision, epoch)
                writer.add_scalar("Recall/val", val_recall, epoch)
                writer.add_scalar("Average Precision/val", val_avg_precision, epoch)

                train_losses.append(train_loss)
                val_losses.append(val_loss)
                train_accuracies.append(train_acc)
                val_accuracies.append(val_acc)

                print(f"Epoch {epoch}, train loss: {train_loss:.4f}, val loss: {val_loss:.4f}")

            # Save and plot the metrics and plots
            plot_loss_accuracy_curve(train_losses, val_losses, train_accuracies, val_accuracies, pdf_pages, title=f"Loss and Accuracy for Fold {idx_fold}")
            plot_roc_curve(val_all_labels, val_all_preds, pdf_pages, title=f"ROC Curve for Fold {idx_fold}")
            plot_precision_recall_curve(val_all_labels, val_all_preds, pdf_pages, title=f"Precision-Recall Curve for Fold {idx_fold}")
            plot_class_distribution(kfold_train_dataset.class_to_idx, pdf_pages, title=f"Class Distribution for Fold {idx_fold}")
            
            # Add content to the PowerPoint
            ppt_title = f"Fold {idx_fold} Training and Validation Results"
            ppt_content = f"Training Loss: {train_losses[-1]:.4f}\nValidation Loss: {val_losses[-1]:.4f}\n" \
                          f"Training Accuracy: {train_accuracies[-1]:.2f}%\nValidation Accuracy: {val_accuracies[-1]:.2f}%"
            add_ppt_slide(ppt, ppt_title, ppt_content)

            # Add plots to PowerPoint (e.g., ROC Curve, Precision-Recall Curve, Loss-Accuracy)
            plot_loss_accuracy_curve(train_losses, val_losses, train_accuracies, val_accuracies, ppt, title=f"Loss & Accuracy for Fold {idx_fold}")
            plot_roc_curve(val_all_labels, val_all_preds, ppt, title=f"ROC Curve for Fold {idx_fold}")
            plot_precision_recall_curve(val_all_labels, val_all_preds, ppt, title=f"Precision-Recall Curve for Fold {idx_fold}")
            plot_class_distribution(kfold_train_dataset.class_to_idx, ppt, title=f"Class Distribution for Fold {idx_fold}")

    # Save PowerPoint presentation
    ppt_output_path = os.path.join(run_output_dir, f"training_report_run_{run}.pptx")
    ppt.save(ppt_output_path)

    print(f"PowerPoint report saved to {ppt_output_path}")








