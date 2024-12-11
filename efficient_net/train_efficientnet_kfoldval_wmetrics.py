from efficientnet_pytorch import EfficientNet
from efficientnet_pytorch import EfficientNet
from efficient_net.options import OptionsModelType
from efficient_net.utils import load_efficient_net_model
from efficient_net.split_dataset_kfoldcrossval import get_kfold_directory_list
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
import numpy as np
import os
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from sklearn.preprocessing import LabelBinarizer
from sklearn.metrics import (
    f1_score,
    precision_score,
    recall_score,
    average_precision_score,
)
import torch
from torch import nn
import torch.optim as optim
from torch.utils.data import DataLoader
from torch.utils.data.dataloader import default_collate
from torch.utils.tensorboard import SummaryWriter
import torchvision.transforms.functional as F
from torchvision import datasets, transforms
from tqdm import tqdm
from utils.plotting import (
    change_matplotlib_font_params,
    plot_class_distribution,
    plot_loss_accuracy_curve,
    plot_roc_curve,
    plot_precision_recall_curve,
)


change_matplotlib_font_params()

# Path and settings
PATH_LOAD_FROM_DISK = None  # "/home/matthias/workspace/Coding/00_vista_medizina/10_weights/efficient_net/2024-12-02_bf_kaggle/efficientnet-b4_BEST.pth"
# "/home/matthias/workspace/Coding/00_vista_medizina/10_weights/efficient_net/2024-12-02_bf_kaggle/efficientnet-b5_epoch30_BEST.pth"
# "/home/matthias/workspace/Coding/00_vista_medizina/10_weights/efficient_net/2024-12-02_bf_kaggle/efficientnet-b4_BEST.pth"
PATH_DATASET = "/home/matthias/workspace/Coding/00_vista_medizina/00_data/2024-11-08/kaggle_bone_fracture_detection_small_2024-11-08/bf_orig_bgr_classification_2cl/kfold"
PATH_OUTPUT_DIR = "/home/matthias/workspace/Coding/00_vista_medizina/vista_bone_frac/efficient_net/output_2cl_5fold_1run_5e-6_30ep"

LR_INITIAL = 5e-6  # 5e-6  # best: 1e-6
LR_PLATEAU_REDUCER = True  # const lr: False
LR_PLAT_FACTOR = 0.5  # FLOAT, pytorch: 0.1
LR_PLAT_PATIENCE = 3  # INT, pytorch:3
MODEL_TYPE = OptionsModelType.B4  #  #
if MODEL_TYPE == OptionsModelType.B5:
    BATCH_SIZE = 10
elif MODEL_TYPE == OptionsModelType.B4:
    BATCH_SIZE = 22
# K_FOLDS = 5  # Number of folds for cross-validation
NUM_EPOCHS = 31  # Number of epochs to train
NUM_RUNS = 1  # 3  # Number of runs for random splits
LIST_SAVE_EPOCHS = [20, 30]
SAVE_CURRENT_BEST_MODEL = False  # False

# Device setup
device = torch.device("cuda" if torch.cuda.is_available() else "cpu")

"""
def plot_roc_curve(
    y_true,
    y_pred_probs,
    save_object,
    title="ROC Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # Convert labels to one-hot format for multiclass classification
    lb = LabelBinarizer()
    y_true_bin = lb.fit_transform(y_true)

    if not isinstance(y_true_bin, np.ndarray):
        y_true_bin = np.array(y_true_bin)
    if not isinstance(y_pred_probs, np.ndarray):
        y_pred_probs = np.array(y_pred_probs)

    # Ensure y_pred_probs is a NumPy array
    y_pred_probs = np.array(y_pred_probs)

    # Check if y_pred_probs is 2D (for multiclass)
    if y_pred_probs.ndim == 1:
        # If it's 1D, expand it to 2D for binary classification (one class vs rest)
        y_pred_probs = np.expand_dims(y_pred_probs, axis=1)

    # Plot ROC curve for each class
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    n_classes = y_true_bin.shape[1]
    fpr, tpr, roc_auc = {}, {}, {}

    for i in range(n_classes):
        fpr[i], tpr[i], _ = roc_curve(y_true_bin[:, i], y_pred_probs[:, i])
        roc_auc[i] = auc(fpr[i], tpr[i])
        plt.plot(fpr[i], tpr[i], label=f"Class {lb.classes_[i]} (AUC = {roc_auc[i]:.2f})")

    plt.plot([0, 1], [0, 1], color="gray", linestyle="--")  # Diagonal line (chance level)
    plt.xlabel("False Positive Rate")
    plt.xlim([0.0, 1.0])
    plt.ylim([0.0, 1.0])
    plt.grid(color="gray", linestyle="--", linewidth=0.5)
    plt.ylabel("True Positive Rate")
    plt.title(title)
    plt.legend(loc="lower right")

    # Save the plot to the provided PDF
    save_object.savefig()
    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_precision_recall_curve(
    y_true,
    y_pred_probs,
    save_object,
    title="Precision-Recall Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # Convert labels to one-hot format for multiclass classification
    lb = LabelBinarizer()
    y_true_bin = lb.fit_transform(y_true)

    if not isinstance(y_true_bin, np.ndarray):
        y_true_bin = np.array(y_true_bin)
    if not isinstance(y_pred_probs, np.ndarray):
        y_pred_probs = np.array(y_pred_probs)

    # Ensure y_pred_probs is a NumPy array
    y_pred_probs = np.array(y_pred_probs)

    # Check if y_pred_probs is 2D (for multiclass)
    if y_pred_probs.ndim == 1:
        # If it's 1D, expand it to 2D for binary classification (one class vs rest)
        y_pred_probs = np.expand_dims(y_pred_probs, axis=1)

    # Plot Precision-Recall curve for each class
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    n_classes = y_true_bin.shape[1]

    for i in range(n_classes):
        precision, recall, _ = precision_recall_curve(y_true_bin[:, i], y_pred_probs[:, i])
        avg_precision = average_precision_score(y_true_bin[:, i], y_pred_probs[:, i])
        plt.plot(recall, precision, label=f"Class {lb.classes_[i]} (AP = {avg_precision:.2f})")

    plt.xlabel("Recall")
    plt.ylabel("Precision")
    plt.title(title)
    plt.xlim([0.0, 1.0])
    plt.ylim([0.0, 1.0])
    plt.grid(color="gray", linestyle="--", linewidth=0.5)
    plt.legend(loc="lower left")

    # Save the plot to the provided PDF
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)

    plt.close()


def plot_loss_curve(
    losses,
    save_object,
    title="Loss Curve",  # tensorboard_writer: SummaryWriter = None
):
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    plt.plot(losses, color="red", lw=2)
    plt.xlabel("Epochs")
    plt.ylabel("Loss")
    plt.title(title)
    save_object.savefig()
    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_class_distribution(
    list_class_occurences,
    save_object,
    title="Class Distribution",
    # tensorboard_writer: SummaryWriter = None,
):
    # sns.barplot(x=list(class_counts.keys()), y=list(class_counts.values()))
    # fixed bin size
    bins = np.arange(0, max(list_class_occurences) + 1, 1)  # fixed bin size

    plt.xlim([min(list_class_occurences), max(list_class_occurences)])
    # fig = plt.hist(list_class_occurences, bins=bins, alpha=0.5)
    plt.hist(list_class_occurences, bins=bins, alpha=0.5)
    plt.xlabel("Classes")
    plt.ylabel("Count")
    plt.title(title)
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_loss_accuracy_curve(
    list_train_losses,
    list_val_losses,
    list_train_accuracies,
    list_val_accuracies,
    save_object,
    titles: list,
    # tensorboard_writer: SummaryWriter = None,
):
    # fig = plt.figure(figsize=(10, 5))
    # plt.figure(figsize=(10, 5))
    plt.subplot(1, 2, 1)
    plt.title(titles[0])
    plt.plot(range(NUM_EPOCHS), list_train_losses, label="Training Loss")
    plt.plot(range(NUM_EPOCHS), list_val_losses, label="Validation Loss")
    plt.xlabel("Epochs")
    plt.ylabel("Loss")
    plt.xlim(0.0)
    plt.ylim(0.0)
    plt.grid(color="gray", linestyle="--", linewidth=0.5)
    plt.legend()

    plt.subplot(1, 2, 2)
    plt.title(titles[1])
    plt.plot(range(NUM_EPOCHS), list_train_accuracies, label="Training Accuracy")
    plt.plot(range(NUM_EPOCHS), list_val_accuracies, label="Validation Accuracy")
    plt.xlabel("Epochs")
    plt.ylabel("Accuracy (%)")
    plt.xlim(0.0)
    plt.ylim([0, 100])
    plt.grid(color="gray", linestyle="--", linewidth=0.5)
    plt.legend(loc="lower right")

    plt.tight_layout()
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


# Helper function to create a PowerPoint presentation slide
def add_ppt_slide(ppt, title, content, image_path=None):
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Use the 'Title and Content' layout

    # Set slide title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Add content (text or image)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(9))  # Add image
    else:
        content_box = slide.shapes.placeholders[1]  # Text box
        content_box.text = content

"""


# Collate function
def custom_collate_fn(batch):
    # Ensure all images in the batch have the same size
    return default_collate(batch)


# Model transforms
class ResizeWithAspectRatioAndPadding:
    def __init__(self, target_size):
        self.target_size = target_size

    def __call__(self, img):
        # Get the original dimensions of the image
        width, height = img.size
        aspect_ratio = width / height

        # Compute the new dimensions while maintaining the aspect ratio
        if aspect_ratio > 1:
            new_width = self.target_size
            new_height = int(self.target_size / aspect_ratio)
        else:
            new_height = self.target_size
            new_width = int(self.target_size * aspect_ratio)

        # Resize the image
        img = F.resize(img, (new_height, new_width))

        # Now, pad the image to match the target size
        padding = (0, 0, self.target_size - new_width, self.target_size - new_height)
        img = F.pad(img, padding, fill=0)
        return img


# Transforms for training and validation
train_transforms = transforms.Compose(
    [
        ResizeWithAspectRatioAndPadding(MODEL_TYPE.get_resolution()),
        transforms.RandomHorizontalFlip(),
        transforms.RandomVerticalFlip(),
        transforms.RandomRotation(20),
        transforms.ToTensor(),
        transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
    ]
)

val_transforms = transforms.Compose(
    [
        ResizeWithAspectRatioAndPadding(MODEL_TYPE.get_resolution()),
        transforms.ToTensor(),
        transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
    ]
)

"""list_paths_kfold_dirs = get_kfold_directory_list(path_image_dir=PATH_DATASET)
# Loop over each fold
list_paths_fold_train_dirs = [None] * len(list_paths_kfold_dirs)
list_paths_fold_val_dirs = [None] * len(list_paths_kfold_dirs)

for dict_path_dir in list_paths_kfold_dirs:
    idx_fold = int(dict_path_dir["idx_fold"])
    list_paths_fold_train_dirs[idx_fold] = dict_path_dir["train"]
    list_paths_fold_val_dirs[idx_fold] = dict_path_dir["val"]"""
dict_kfold_dirs = get_kfold_directory_list(path_image_dir=PATH_DATASET)
list_paths_fold_train_dirs = dict_kfold_dirs["list_paths_fold_train_dirs"]
list_paths_fold_val_dirs = dict_kfold_dirs["list_paths_fold_val_dirs"]

# list_paths_fold_train_dirs = [list_paths_fold_train_dirs[0]]
# list_paths_fold_val_dirs = [list_paths_fold_val_dirs[0]]

# Start training process
for run in range(NUM_RUNS):
    print("++++++++++++++++++++++++++++++++++++++++++++++++++++")
    print(f"\nTRAIN RUN NUMBER: {run}/{NUM_RUNS-1}")

    # Create a directory for the current run
    run_output_dir = os.path.join(PATH_OUTPUT_DIR, f"run_{run}")
    os.makedirs(run_output_dir, exist_ok=True)
    writer = SummaryWriter(os.path.join(run_output_dir, "tensorboard"))

    # Create a PowerPoint presentation object
    # ppt = Presentation()

    with PdfPages(os.path.join(run_output_dir, f"training_report_run_{run}.pdf")) as pdf_pages:

        for idx_fold, (path_kfold_train_dir, path_kfold_val_dir) in enumerate(
            zip(list_paths_fold_train_dirs, list_paths_fold_val_dirs)
        ):

            print(f"\n----Fold {idx_fold}----------------------------")
            print("--Creating fold output directory...")

            fold_output_dir = os.path.join(run_output_dir, f"fold_{idx_fold}")
            os.makedirs(fold_output_dir, exist_ok=True)

            print("--Creating datasets...")
            kfold_train_dataset = datasets.ImageFolder(
                path_kfold_train_dir, transform=train_transforms
            )
            kfold_val_dataset = datasets.ImageFolder(path_kfold_val_dir, transform=val_transforms)
            print("--Creating dataloaders...")
            train_fold_loader = DataLoader(
                kfold_train_dataset,
                batch_size=BATCH_SIZE,
                shuffle=True,
                num_workers=4,
                collate_fn=custom_collate_fn,
            )
            val_fold_loader = DataLoader(
                kfold_val_dataset,
                batch_size=BATCH_SIZE,
                shuffle=False,
                num_workers=4,
                collate_fn=custom_collate_fn,
            )

            print("--Loading the model...")
            if PATH_LOAD_FROM_DISK is not None:
                model = load_efficient_net_model(
                    path_model=PATH_LOAD_FROM_DISK,
                    model_type=MODEL_TYPE,
                    num_classes=len(kfold_train_dataset.classes),
                )

            else:
                model = EfficientNet.from_pretrained(
                    MODEL_TYPE, num_classes=len(kfold_train_dataset.classes)
                )
            model.train()
            model.to(device)

            criterion = nn.CrossEntropyLoss()  # CrossEntropyLoss for multi-class classification
            optimizer = optim.Adam(model.parameters(), lr=LR_INITIAL)  # Adam optimizer

            scaler = torch.amp.GradScaler()
            if LR_PLATEAU_REDUCER:
                scheduler = torch.optim.lr_scheduler.ReduceLROnPlateau(
                    optimizer,
                    mode="min",
                    factor=LR_PLAT_FACTOR,
                    patience=LR_PLAT_PATIENCE,
                    verbose=True,
                )

            list_train_losses, list_val_losses = [], []
            list_train_accuracies, list_val_accuracies = [], []
            list_train_f1 = []
            list_train_precision = []
            list_train_recall = []
            list_train_average_precision = []

            list_val_f1 = []
            list_val_precision = []
            list_val_recall = []
            list_val_average_precision = []

            highest_val_acc = 0.0
            print("--Starting the training process...")
            for epoch in range(NUM_EPOCHS):
                print(
                    f"\n+++++ Run:{run}, fold:{idx_fold}, epoch: {epoch}/{NUM_EPOCHS-1} ++++++++++++++++++++++++++++++++++++++++++++++"
                )
                current_lr = optimizer.param_groups[0]["lr"]
                print(f"--learning rate: {current_lr:.1e}")

                # Training Phase
                model.train()
                train_running_loss = 0.0
                train_correct_preds = 0
                train_total_preds = 0
                train_all_preds = []
                train_all_labels = []
                train_all_pred_probs = (
                    []
                )  # To store predicted probabilities for average_precision_score

                print("--training loop")
                for images, labels in tqdm(train_fold_loader):
                    images, labels = images.to(device), labels.to(device)

                    optimizer.zero_grad()

                    with torch.amp.autocast(device_type=str(device)):  # Mixed precision training
                        output = model(images)
                        loss = criterion(output, labels)

                    scaler.scale(loss).backward()
                    scaler.step(optimizer)
                    scaler.update()

                    train_running_loss += loss.item()
                    _, predicted = torch.max(output, 1)
                    train_correct_preds += (predicted == labels).sum().item()
                    train_total_preds += labels.size(0)

                    train_all_preds.extend(predicted.cpu().numpy())
                    train_all_labels.extend(labels.cpu().numpy())

                    # Calculate training metrics
                    train_loss = train_running_loss / len(train_fold_loader)
                    train_acc = 100 * train_correct_preds / train_total_preds

                    # Collect predicted probabilities for average_precision_score
                    pred_probs = torch.softmax(
                        output, dim=1
                    )  # Apply softmax to get probabilities for each class
                    train_all_pred_probs.extend(pred_probs.cpu().detach().numpy())

                # Calculate F1, Precision, Recall, and Average Precision using the correct inputs
                train_f1 = f1_score(train_all_labels, train_all_preds, average="weighted")
                train_precision = precision_score(
                    train_all_labels, train_all_preds, average="weighted"
                )
                train_recall = recall_score(train_all_labels, train_all_preds, average="weighted")
                if len(kfold_train_dataset.classes) > 2:
                    train_average_precision = average_precision_score(
                        LabelBinarizer().fit_transform(
                            train_all_labels
                        ),  # Convert true labels to one-hot format
                        train_all_pred_probs,  # Predicted probabilities for each class
                        average="weighted",
                    )
                else:
                    train_average_precision = train_precision

                # Log to TensorBoard
                writer.add_scalar("Loss/train", train_loss, epoch)
                writer.add_scalar("Accuracy/train", train_acc, epoch)
                writer.add_scalar("F1/train", train_f1, epoch)
                writer.add_scalar("Precision/train", train_precision, epoch)
                writer.add_scalar("Recall/train", train_recall, epoch)
                writer.add_scalar("Average Precision/train", train_average_precision, epoch)

                # Validation
                model.eval()
                val_running_loss, val_correct_preds, val_total_preds = 0.0, 0, 0
                print("--validation loop")
                val_all_preds, val_all_labels = [], []
                val_all_pred_probs = (
                    []
                )  # To store predicted probabilities for average_precision_score

                with torch.no_grad():
                    for images, labels in tqdm(val_fold_loader):
                        images, labels = images.to(device), labels.to(device)
                        outputs = model(images)
                        loss = criterion(outputs, labels)
                        val_running_loss += loss.item()
                        _, predicted = torch.max(outputs, 1)
                        val_correct_preds += (predicted == labels).sum().item()
                        val_total_preds += labels.size(0)

                        val_all_preds.extend(predicted.cpu().numpy())
                        val_all_labels.extend(labels.cpu().numpy())

                        # Collect predicted probabilities for average_precision_score
                        pred_probs = torch.softmax(
                            outputs, dim=1
                        )  # Apply softmax to get probabilities for each class
                        val_all_pred_probs.extend(pred_probs.cpu().detach().numpy())

                val_loss = val_running_loss / len(val_fold_loader)
                val_acc = 100 * val_correct_preds / val_total_preds
                val_f1 = f1_score(val_all_labels, val_all_preds, average="weighted")
                val_precision = precision_score(val_all_labels, val_all_preds, average="weighted")
                val_recall = recall_score(val_all_labels, val_all_preds, average="weighted")
                if len(kfold_val_dataset.classes) > 2:
                    val_average_precision = average_precision_score(
                        LabelBinarizer().fit_transform(
                            val_all_labels
                        ),  # Convert true labels to one-hot format
                        val_all_pred_probs,  # Predicted probabilities for each class
                        average="weighted",
                    )
                else:
                    val_average_precision = val_precision

                # Log validation metrics to TensorBoard
                writer.add_scalar("Loss/val", val_loss, epoch)
                writer.add_scalar("Accuracy/val", val_acc, epoch)
                writer.add_scalar("F1/val", val_f1, epoch)
                writer.add_scalar("Precision/val", val_precision, epoch)
                writer.add_scalar("Recall/val", val_recall, epoch)
                writer.add_scalar("Average Precision/val", val_average_precision, epoch)
                if LR_PLATEAU_REDUCER:
                    scheduler.step(val_loss)
                    print(
                        f"--LR scheduler: ReduceLROnPlateau reducing learning rate to {scheduler.get_last_lr()[0]:.1e}"
                    )

                list_train_losses.append(train_loss)
                list_val_losses.append(val_loss)

                list_train_f1.append(train_f1)
                list_train_precision.append(train_precision)
                list_train_recall.append(train_recall)
                list_train_average_precision.append(train_average_precision)

                list_val_f1.append(val_f1)
                list_val_precision.append(val_precision)
                list_val_recall.append(val_recall)
                list_val_average_precision.append(val_average_precision)

                list_train_accuracies.append(train_acc)
                list_val_accuracies.append(val_acc)

                # Print out training and validation losses and accuracies
                print(
                    f"---- train loss: {train_loss:.4f}, train accuracy: {train_acc:.2f}%, "
                    f"---- validation loss: {val_loss:.4f}, validation accuracy: {val_acc:.2f}%"
                )

                # Save model with highest validation accuracy
                if (
                    SAVE_CURRENT_BEST_MODEL
                    and epoch > int(NUM_EPOCHS * 0.5)
                    and val_acc > highest_val_acc
                ):
                    highest_val_acc = val_acc
                    print("==> Model saved with highest validation accuracy!")
                    torch.save(
                        model.state_dict(),
                        os.path.join(
                            fold_output_dir,
                            str(MODEL_TYPE.value) + "_epoch" + str(epoch) + "_BEST.pth",
                        ),
                    )

                # Save the model at specified epochs
                if (
                    LIST_SAVE_EPOCHS is not None
                    and len(LIST_SAVE_EPOCHS) > 0
                    and epoch in LIST_SAVE_EPOCHS
                ):
                    print("==> Saving model to disk...")
                    torch.save(
                        model.state_dict(),
                        os.path.join(
                            fold_output_dir,
                            str(MODEL_TYPE.value) + "_epoch" + str(epoch) + "_SCHED.pth",
                        ),
                    )
                # Free up memory after each epoch
                del images, labels, output, loss
                torch.cuda.empty_cache()

            print(
                "--finished computing fold number "
                + str(idx_fold)
                + ", run number "
                + str(run)
                + "."
            )
            """
            # Plot and save the loss curves
            plt.figure(figsize=(10, 5))
            plt.subplot(1, 2, 1)
            plt.plot(range(NUM_EPOCHS), list_train_losses, label="Training Loss")
            plt.plot(range(NUM_EPOCHS), list_val_losses, label="Validation Loss")
            plt.xlabel("Epochs")
            plt.ylabel("Loss")
            plt.legend()

            plt.subplot(1, 2, 2)
            plt.plot(range(NUM_EPOCHS), list_train_accuracies, label="Training Accuracy")
            plt.plot(range(NUM_EPOCHS), list_val_accuracies, label="Validation Accuracy")
            plt.xlabel("Epochs")
            plt.ylabel("Accuracy (%)")
            plt.legend()

            plt.tight_layout()
            plt.savefig(os.path.join(fold_output_dir, "loss_accuracy_curve.png"))
            plt.close()
            print("--finished saving loss plots.")

            writer.flush() 
            print(f"Finished fold {idx_fold} of run {run + 1}/{NUM_RUNS}.")
            writer.close()
            """

            def plot_metrics(save_object, title: str):
                # Plot and save the metrics
                # plt.figure(figsize=(10, 5))
                plt.subplot(2, 2, 1)
                plt.title(title)
                plt.plot(range(NUM_EPOCHS), list_train_f1, label="Train F1 score")
                plt.plot(range(NUM_EPOCHS), list_val_f1, label="Val F1 score")
                plt.xlabel("Epochs")
                plt.ylabel("F1 Score")
                plt.xlim(0)
                plt.ylim([0.0, 1.0])
                plt.grid(color="gray", linestyle="--", linewidth=0.5)
                plt.legend(loc="lower right")

                plt.subplot(2, 2, 2)
                plt.plot(range(NUM_EPOCHS), list_train_precision, label="Train precision")
                plt.plot(range(NUM_EPOCHS), list_val_precision, label="Val precision")
                plt.xlabel("Epochs")
                plt.ylabel("Precision")
                plt.xlim(0)
                plt.ylim([0.0, 1.0])
                plt.grid(color="gray", linestyle="--", linewidth=0.5)
                plt.legend(loc="lower right")

                plt.subplot(2, 2, 3)
                plt.plot(range(NUM_EPOCHS), list_train_recall, label="Train recall")
                plt.plot(range(NUM_EPOCHS), list_val_recall, label="Val recall")
                plt.xlabel("Epochs")
                plt.ylabel("Recall")
                plt.xlim(0)
                plt.ylim([0.0, 1.0])
                plt.grid(color="gray", linestyle="--", linewidth=0.5)
                plt.legend(loc="lower right")

                plt.subplot(2, 2, 4)
                plt.plot(
                    range(NUM_EPOCHS), list_train_average_precision, label="Train avg. precision"
                )
                plt.plot(range(NUM_EPOCHS), list_val_average_precision, label="Val avg. precision")
                plt.xlabel("Epochs")
                plt.ylabel("Average Precision")
                plt.ylim([0.0, 1.0])
                plt.xlim(0)
                plt.grid(color="gray", linestyle="--", linewidth=0.5)
                plt.legend(loc="lower right")

                plt.tight_layout()
                # Save the plot to the provided PDF
                save_object.savefig()
                plt.close()

            def print_all_plots_to_pdf(save_object):
                # Save and plot the metrics and plots
                plot_class_distribution(
                    kfold_train_dataset.__dict__["targets"],
                    save_object,
                    title=f"Training Class Distribution for Fold {idx_fold}",
                    # tensorboard_writer=writer,
                )
                plot_class_distribution(
                    kfold_val_dataset.__dict__["targets"],
                    save_object,
                    title=f"Validation Class Distribution for Fold {idx_fold}",
                    # tensorboard_writer=writer,
                )
                plot_loss_accuracy_curve(
                    NUM_EPOCHS,
                    list_train_losses,
                    list_val_losses,
                    list_train_accuracies,
                    list_val_accuracies,
                    save_object,
                    titles=[f"Loss, fold: {idx_fold}", f"Accuracy, fold: {idx_fold}"],
                    # tensorboard_writer=writer,
                )
                plot_roc_curve(
                    train_all_labels,  # True labels
                    train_all_pred_probs,  # Predicted probabilities (2D array)
                    save_object,  # The PDF object where the plot will be saved
                    title=f"Training ROC Curve for Fold {idx_fold}",  # Title for the plot
                    # tensorboard_writer=writer,
                )
                plot_roc_curve(
                    val_all_labels,  # True labels
                    val_all_pred_probs,  # Predicted probabilities (2D array)
                    save_object,  # The PDF object where the plot will be saved
                    title=f"Validation ROC Curve for Fold {idx_fold}",  # Title for the plot
                    # tensorboard_writer=writer,
                )
                plot_precision_recall_curve(
                    train_all_labels,  # True labels
                    train_all_pred_probs,  # Predicted probabilities (2D array)
                    save_object,  # The PDF object where the plot will be saved
                    title=f"Training Precision-Recall Curve for Fold {idx_fold}",  # Title for the plot
                    # tensorboard_writer=writer,
                )
                plot_precision_recall_curve(
                    val_all_labels,  # True labels
                    val_all_pred_probs,  # Predicted probabilities (2D array)
                    save_object,  # The PDF object where the plot will be saved
                    title=f"Validation Precision-Recall Curve for Fold {idx_fold}",  # Title for the plot
                    # tensorboard_writer=writer,
                )
                plot_metrics(
                    save_object,
                    title=f"Training and validation metrics for fold: {idx_fold}",
                    # tensorboard_writer=writer,
                )

            print_all_plots_to_pdf(pdf_pages)

            writer.flush()
            writer.close()
            """
            # Add content to the PowerPoint
            ppt_title = f"Fold {idx_fold} Training and Validation Results"
            ppt_content = (
                f"Training Loss: {list_train_losses[-1]:.4f}\nValidation Loss: {list_val_losses[-1]:.4f}\n"
                f"Training Accuracy: {list_train_accuracies[-1]:.2f}%\nValidation Accuracy: {list_val_accuracies[-1]:.2f}%"
            )
            add_ppt_slide(ppt, ppt_title, ppt_content)
            print_all_plots_to_object(ppt)
            """
            """
            # Add plots to PowerPoint (e.g., ROC Curve, Precision-Recall Curve, Loss-Accuracy)
            plot_loss_accuracy_curve(
                list_train_losses,
                list_val_losses,
                list_train_accuracies,
                list_val_accuracies,
                ppt,
                title=f"Loss & Accuracy for Fold {idx_fold}",
            )
            plot_roc_curve(
                val_all_labels,  # True labels
                val_all_pred_probs,  # Predicted probabilities (2D array)
                ppt,  # The PDF object where the plot will be saved
                title=f"ROC Curve for Fold {idx_fold}",  # Title for the plot
            )
            plot_precision_recall_curve(
                val_all_labels,  # True labels
                val_all_pred_probs,  # Predicted probabilities (2D array)
                ppt,  # The PDF object where the plot will be saved
                title=f"Precision-Recall Curve for Fold {idx_fold}",  # Title for the plot
            )
            plot_class_distribution(
                kfold_train_dataset.class_to_idx, ppt, title=f"Class Distribution for Fold {idx_fold}"
            )
            """

    # Save PowerPoint presentation
    # ppt_output_path = os.path.join(run_output_dir, f"training_report_run_{run}.pptx")
    # ppt.save(ppt_output_path)

    # print(f"PowerPoint report saved to {ppt_output_path}")
    print(f"Finished run {run + 1}/{NUM_RUNS}.")
