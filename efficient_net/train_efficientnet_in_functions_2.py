from efficientnet_pytorch import EfficientNet
from efficientnet_pytorch import EfficientNet
from efficient_net.options import OptionsModelType
from efficient_net.utils import load_efficient_net_model
from efficient_net.split_dataset_kfoldcrossval import get_kfold_directory_list
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
import numpy as np
import os
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
from sklearn.preprocessing import LabelBinarizer
from sklearn.metrics import (
    f1_score,
    precision_score,
    precision_recall_curve,
    recall_score,
    average_precision_score,
    roc_auc_score,
    roc_curve,
    auc,
)
import torch
from torch import nn
import torch.optim as optim
from torch.utils.data import DataLoader
from torch.utils.data.dataloader import default_collate
from torch.utils.tensorboard import SummaryWriter
import torchvision.transforms.functional as F
from torchvision import datasets, transforms
from tqdm import tqdm


# Model transforms
class ResizeWithAspectRatioAndPadding:
    def __init__(self, target_size):
        self.target_size = target_size

    def __call__(self, img):
        # Get the original dimensions of the image
        width, height = img.size
        aspect_ratio = width / height

        # Compute the new dimensions while maintaining the aspect ratio
        if aspect_ratio > 1:
            new_width = self.target_size
            new_height = int(self.target_size / aspect_ratio)
        else:
            new_height = self.target_size
            new_width = int(self.target_size * aspect_ratio)

        # Resize the image
        img = F.resize(img, (new_height, new_width))

        # Now, pad the image to match the target size
        padding = (0, 0, self.target_size - new_width, self.target_size - new_height)
        img = F.pad(img, padding, fill=0)
        return img


# Path and settings
PATH_LOAD_FROM_DISK = "/home/matthias/workspace/Coding/00_vista_medizina/10_weights/efficient_net/2024-12-02_bf_kaggle/efficientnet-b4_BEST.pth"
# "/home/matthias/workspace/Coding/00_vista_medizina/10_weights/efficient_net/2024-12-02_bf_kaggle/efficientnet-b5_epoch30_BEST.pth"
# "/home/matthias/workspace/Coding/00_vista_medizina/10_weights/efficient_net/2024-12-02_bf_kaggle/efficientnet-b4_BEST.pth"
PATH_DATASET = "/home/matthias/workspace/Coding/00_vista_medizina/00_data/2024-11-08/kaggle_bone_fracture_detection_small_2024-11-08/bf_orig_bgr_classification_7cl_kfold/w_folds"
PATH_OUTPUT_DIR = "/home/matthias/workspace/Coding/00_vista_medizina/vista_bone_frac/efficient_net/output_7cl_5fold_3runs"

LR_INITIAL = 5e-6  # best: 1e-6
LR_PLATEAU_REDUCER = True  # const lr: False
LR_PLAT_FACTOR = 0.5  # FLOAT, pytorch: 0.1
LR_PLAT_PATIENCE = 5  # INT, pytorch:3
MODEL_TYPE = OptionsModelType.B4  #  #
if MODEL_TYPE == OptionsModelType.B5:
    BATCH_SIZE = 10
elif MODEL_TYPE == OptionsModelType.B4:
    BATCH_SIZE = 22
# K_FOLDS = 5  # Number of folds for cross-validation
NUM_EPOCHS = 31  # 31  # Number of epochs to train
NUM_RUNS = 1  # 3  # Number of runs for random splits
LIST_SAVE_EPOCHS = [10, 20, 30]
SAVE_CURRENT_BEST_MODEL = True  # False

# Device setup
device = torch.device("cuda" if torch.cuda.is_available() else "cpu")


# Transforms for training and validation
train_transforms = transforms.Compose(
    [
        ResizeWithAspectRatioAndPadding(MODEL_TYPE.get_resolution()),
        transforms.RandomHorizontalFlip(),
        transforms.RandomVerticalFlip(),
        transforms.RandomRotation(20),
        transforms.ToTensor(),
        transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
    ]
)

val_transforms = transforms.Compose(
    [
        ResizeWithAspectRatioAndPadding(MODEL_TYPE.get_resolution()),
        transforms.ToTensor(),
        transforms.Normalize(mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225]),
    ]
)


def plot_roc_curve(
    y_true,
    y_pred_probs,
    save_object,
    title="ROC Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # Convert labels to one-hot format for multiclass classification
    lb = LabelBinarizer()
    y_true_bin = lb.fit_transform(y_true)

    if not isinstance(y_true_bin, np.ndarray):
        y_true_bin = np.array(y_true_bin)
    if not isinstance(y_pred_probs, np.ndarray):
        y_pred_probs = np.array(y_pred_probs)

    # Ensure y_pred_probs is a NumPy array
    y_pred_probs = np.array(y_pred_probs)

    # Check if y_pred_probs is 2D (for multiclass)
    if y_pred_probs.ndim == 1:
        # If it's 1D, expand it to 2D for binary classification (one class vs rest)
        y_pred_probs = np.expand_dims(y_pred_probs, axis=1)

    # Plot ROC curve for each class
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    n_classes = y_true_bin.shape[1]
    fpr, tpr, roc_auc = {}, {}, {}

    for i in range(n_classes):
        fpr[i], tpr[i], _ = roc_curve(y_true_bin[:, i], y_pred_probs[:, i])
        roc_auc[i] = auc(fpr[i], tpr[i])
        plt.plot(fpr[i], tpr[i], label=f"Class {lb.classes_[i]} (AUC = {roc_auc[i]:.2f})")

    plt.plot([0, 1], [0, 1], color="gray", linestyle="--")  # Diagonal line (chance level)
    plt.xlabel("False Positive Rate")
    plt.ylabel("True Positive Rate")
    plt.title(title)
    plt.legend(loc="lower right")

    # Save the plot to the provided PDF
    save_object.savefig()
    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_precision_recall_curve(
    y_true,
    y_pred_probs,
    save_object,
    title="Precision-Recall Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # Convert labels to one-hot format for multiclass classification
    lb = LabelBinarizer()
    y_true_bin = lb.fit_transform(y_true)

    if not isinstance(y_true_bin, np.ndarray):
        y_true_bin = np.array(y_true_bin)
    if not isinstance(y_pred_probs, np.ndarray):
        y_pred_probs = np.array(y_pred_probs)

    # Ensure y_pred_probs is a NumPy array
    y_pred_probs = np.array(y_pred_probs)

    # Check if y_pred_probs is 2D (for multiclass)
    if y_pred_probs.ndim == 1:
        # If it's 1D, expand it to 2D for binary classification (one class vs rest)
        y_pred_probs = np.expand_dims(y_pred_probs, axis=1)

    # Plot Precision-Recall curve for each class
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    n_classes = y_true_bin.shape[1]

    for i in range(n_classes):
        precision, recall, _ = precision_recall_curve(y_true_bin[:, i], y_pred_probs[:, i])
        avg_precision = average_precision_score(y_true_bin[:, i], y_pred_probs[:, i])
        plt.plot(recall, precision, label=f"Class {lb.classes_[i]} (AP = {avg_precision:.2f})")

    plt.xlabel("Recall")
    plt.ylabel("Precision")
    plt.title(title)
    plt.legend(loc="lower left")

    # Save the plot to the provided PDF
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)

    plt.close()


def plot_loss_curve(
    losses,
    save_object,
    title="Loss Curve",  # tensorboard_writer: SummaryWriter = None
):
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    plt.plot(losses, color="red", lw=2)
    plt.xlabel("Epochs")
    plt.ylabel("Loss")
    plt.title(title)
    save_object.savefig()
    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_class_distribution(
    list_class_occurences,
    save_object,
    title="Class Distribution",
    # tensorboard_writer: SummaryWriter = None,
):
    # sns.barplot(x=list(class_counts.keys()), y=list(class_counts.values()))
    # fixed bin size
    bins = np.arange(0, max(list_class_occurences) + 1, 1)  # fixed bin size

    plt.xlim([min(list_class_occurences), max(list_class_occurences)])
    # fig = plt.hist(list_class_occurences, bins=bins, alpha=0.5)
    plt.hist(list_class_occurences, bins=bins, alpha=0.5)
    plt.xlabel("Classes")
    plt.ylabel("Count")
    plt.title(title)
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_loss_accuracy_curve(
    list_train_losses,
    list_val_losses,
    list_train_accuracies,
    list_val_accuracies,
    save_object,
    title="Loss & Accuracy Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # fig = plt.figure(figsize=(10, 5))
    plt.figure(figsize=(10, 5))
    plt.title(title)
    plt.subplot(1, 2, 1)
    plt.plot(range(NUM_EPOCHS), list_train_losses, label="Training Loss")
    plt.plot(range(NUM_EPOCHS), list_val_losses, label="Validation Loss")
    plt.xlabel("Epochs")
    plt.ylabel("Loss")
    plt.legend()

    plt.subplot(1, 2, 2)
    plt.plot(range(NUM_EPOCHS), list_train_accuracies, label="Training Accuracy")
    plt.plot(range(NUM_EPOCHS), list_val_accuracies, label="Validation Accuracy")
    plt.xlabel("Epochs")
    plt.ylabel("Accuracy (%)")
    plt.legend()

    plt.tight_layout()
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


# Helper function to create a PowerPoint presentation slide
def add_ppt_slide(ppt, title, content, image_path=None):
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Use the 'Title and Content' layout

    # Set slide title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Add content (text or image)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(9))  # Add image
    else:
        content_box = slide.shapes.placeholders[1]  # Text box
        content_box.text = content


# Collate function
def custom_collate_fn(batch):
    # Ensure all images in the batch have the same size
    return default_collate(batch)


dict_kfold_dirs = get_kfold_directory_list(path_image_dir=PATH_DATASET)
list_paths_fold_train_dirs = dict_kfold_dirs["list_paths_fold_train_dirs"]
list_paths_fold_val_dirs = dict_kfold_dirs["list_paths_fold_val_dirs"]


def create_directories(run, idx_fold):
    """
    Creates the necessary directories for saving the model, logs, and evaluation results for each run and fold.

    Args:
        run (int): The current run number.
        idx_fold (int): The current fold number in the k-fold cross-validation.

    Returns:
        tuple: The paths to the run and fold directories.
    """
    run_output_dir = os.path.join(PATH_OUTPUT_DIR, f"run_{run}")
    os.makedirs(run_output_dir, exist_ok=True)
    fold_output_dir = os.path.join(run_output_dir, f"fold_{idx_fold}")
    os.makedirs(fold_output_dir, exist_ok=True)
    return run_output_dir, fold_output_dir


def create_dataloaders(path_kfold_train_dir, path_kfold_val_dir):
    """
    Creates data loaders for both the training and validation datasets for a specific fold.

    Args:
        path_kfold_train_dir (str): Path to the training dataset for the current fold.
        path_kfold_val_dir (str): Path to the validation dataset for the current fold.

    Returns:
        tuple: Data loaders for training and validation datasets, along with the raw datasets.
    """
    kfold_train_dataset = datasets.ImageFolder(path_kfold_train_dir, transform=train_transforms)
    kfold_val_dataset = datasets.ImageFolder(path_kfold_val_dir, transform=val_transforms)

    train_fold_loader = DataLoader(
        kfold_train_dataset,
        batch_size=BATCH_SIZE,
        shuffle=True,
        num_workers=4,
        collate_fn=custom_collate_fn,
    )
    val_fold_loader = DataLoader(
        kfold_val_dataset,
        batch_size=BATCH_SIZE,
        shuffle=False,
        num_workers=4,
        collate_fn=custom_collate_fn,
    )

    return train_fold_loader, val_fold_loader, kfold_train_dataset, kfold_val_dataset


def load_model(kfold_train_dataset, path_model=None):
    """
    Loads a pre-trained model or initializes a new model from scratch, based on whether a model path is provided.

    Args:
        kfold_train_dataset (Dataset): The training dataset used to determine the number of output classes.
        path_model (str, optional): The path to a saved model file. If None, a model is initialized from scratch.

    Returns:
        nn.Module: The loaded or newly initialized model.
    """
    if path_model is not None:
        model = load_efficient_net_model(path_model, MODEL_TYPE, len(kfold_train_dataset.classes))
    else:
        model = EfficientNet.from_pretrained(
            MODEL_TYPE, num_classes=len(kfold_train_dataset.classes)
        )

    return model


def train_single_epoch(model, train_fold_loader, optimizer, criterion, scaler, device):
    """
    Executes one epoch of training, including forward pass, loss calculation, backward pass, and parameter update.

    Args:
        model (nn.Module): The model to be trained.
        train_fold_loader (DataLoader): Data loader for the training dataset.
        optimizer (torch.optim.Optimizer): The optimizer used for updating the model's parameters.
        criterion (nn.Module): The loss function used for training.
        scaler (torch.cuda.amp.GradScaler): The gradient scaler for mixed precision training.
        device (torch.device): The device to run the training on (CPU or GPU).

    Returns:
        tuple: The training loss, number of correct predictions, total predictions,
               predicted labels, true labels, and predicted probabilities.
    """
    model.train()
    train_running_loss = 0.0
    train_correct_preds = 0
    train_total_preds = 0
    train_all_preds = []
    train_all_labels = []
    train_all_pred_probs = []

    for images, labels in tqdm(train_fold_loader):
        images, labels = images.to(device), labels.to(device)
        optimizer.zero_grad()

        with torch.amp.autocast(device_type=str(device)):  # Mixed precision training
            output = model(images)
            loss = criterion(output, labels)

        scaler.scale(loss).backward()
        scaler.step(optimizer)
        scaler.update()

        train_running_loss += loss.item()
        _, predicted = torch.max(output, 1)
        train_correct_preds += (predicted == labels).sum().item()
        train_total_preds += labels.size(0)

        train_all_preds.extend(predicted.cpu().numpy())
        train_all_labels.extend(labels.cpu().numpy())

        # Collect predicted probabilities for average_precision_score
        pred_probs = torch.softmax(
            output, dim=1
        )  # Apply softmax to get probabilities for each class
        train_all_pred_probs.extend(pred_probs.cpu().detach().numpy())

    return (
        train_running_loss,
        train_correct_preds,
        train_total_preds,
        train_all_preds,
        train_all_labels,
        train_all_pred_probs,
    )


def validate_single_epoch(model, val_fold_loader, criterion, device):
    """
    Executes one epoch of validation, including forward pass, loss calculation, and performance metrics evaluation.

    Args:
        model (nn.Module): The model to be validated.
        val_fold_loader (DataLoader): Data loader for the validation dataset.
        criterion (nn.Module): The loss function used for validation.
        device (torch.device): The device to run the validation on (CPU or GPU).

    Returns:
        tuple: The validation loss, number of correct predictions, total predictions,
               predicted labels, true labels, and predicted probabilities.
    """
    model.eval()
    val_running_loss, val_correct_preds, val_total_preds = 0.0, 0, 0
    val_all_preds, val_all_labels = [], []
    val_all_pred_probs = []

    with torch.no_grad():
        for images, labels in tqdm(val_fold_loader):
            images, labels = images.to(device), labels.to(device)
            outputs = model(images)
            loss = criterion(outputs, labels)
            val_running_loss += loss.item()
            _, predicted = torch.max(outputs, 1)
            val_correct_preds += (predicted == labels).sum().item()
            val_total_preds += labels.size(0)

            val_all_preds.extend(predicted.cpu().numpy())
            val_all_labels.extend(labels.cpu().numpy())

            # Collect predicted probabilities for average_precision_score
            pred_probs = torch.softmax(
                outputs, dim=1
            )  # Apply softmax to get probabilities for each class
            val_all_pred_probs.extend(pred_probs.cpu().detach().numpy())

    return (
        val_running_loss,
        val_correct_preds,
        val_total_preds,
        val_all_preds,
        val_all_labels,
        val_all_pred_probs,
    )


def compute_metrics(all_labels, all_preds, all_pred_probs):
    """
    Computes performance metrics including F1 score, Precision, Recall, and Average Precision.

    Args:
        all_labels (list): List of true labels.
        all_preds (list): List of predicted labels.
        all_pred_probs (list): List of predicted probabilities for each class.

    Returns:
        tuple: The calculated F1 score, Precision, Recall, and Average Precision.
    """
    f1 = f1_score(all_labels, all_preds, average="weighted")
    precision = precision_score(all_labels, all_preds, average="weighted")
    recall = recall_score(all_labels, all_preds, average="weighted")
    average_precision = average_precision_score(
        LabelBinarizer().fit_transform(all_labels), all_pred_probs, average="weighted"
    )

    return f1, precision, recall, average_precision


def save_model(model, fold_output_dir, epoch, val_acc, highest_val_acc):
    """
    Saves the model if it achieves the highest validation accuracy or at scheduled epochs.

    Args:
        model (nn.Module): The model to be saved.
        fold_output_dir (str): The directory to save the model.
        epoch (int): The current epoch number.
        val_acc (float): The validation accuracy for the current epoch.
        highest_val_acc (float): The highest validation accuracy encountered so far.
    """
    if SAVE_CURRENT_BEST_MODEL and epoch > int(NUM_EPOCHS * 0.5) and val_acc > highest_val_acc:
        highest_val_acc = val_acc
        print("==> Model saved with highest validation accuracy!")
        torch.save(
            model.state_dict(), os.path.join(fold_output_dir, f"{MODEL_TYPE}_epoch{epoch}_BEST.pth")
        )

    # Save the model at specified epochs
    if LIST_SAVE_EPOCHS is not None and len(LIST_SAVE_EPOCHS) > 0 and epoch in LIST_SAVE_EPOCHS:
        print("==> Saving model as scheduled to disk...")
        torch.save(
            model.state_dict(),
            os.path.join(fold_output_dir, f"{MODEL_TYPE}_epoch{epoch}_SCHED.pth"),
        )


def plot_train_and_val_metrics(
    save_object,
    list_train_f1,
    list_val_f1,
    list_train_precision,
    list_val_precision,
    list_train_recall,
    list_val_recall,
    list_train_average_precision,
    list_val_average_precision,
):
    """
    Plots and saves the training and validation metrics (F1, Precision, Recall, and Average Precision) for each epoch.

    Args:
        save_object (PdfPages): The object to save the plots to a PDF.
        list_train_f1 (list): List of F1 scores for training.
        list_val_f1 (list): List of F1 scores for validation.
        list_train_precision (list): List of Precision scores for training.
        list_val_precision (list): List of Precision scores for validation.
        list_train_recall (list): List of Recall scores for training.
        list_val_recall (list): List of Recall scores for validation.
        list_train_average_precision (list): List of Average Precision scores for training.
        list_val_average_precision (list): List of Average Precision scores for validation.
    """
    plt.figure(figsize=(10, 5))
    plt.title("Training and Validation Metrics")
    plt.subplot(2, 2, 1)
    plt.plot(range(NUM_EPOCHS), list_train_f1, label="Training F1 Score")
    plt.plot(range(NUM_EPOCHS), list_val_f1, label="Validation F1 Score")
    plt.xlabel("Epochs")
    plt.ylabel("F1 Score")
    plt.legend()

    plt.subplot(2, 2, 2)
    plt.plot(range(NUM_EPOCHS), list_train_precision, label="Training Precision")
    plt.plot(range(NUM_EPOCHS), list_val_precision, label="Validation Precision")
    plt.xlabel("Epochs")
    plt.ylabel("Precision")
    plt.legend()

    plt.subplot(2, 2, 3)
    plt.plot(range(NUM_EPOCHS), list_train_recall, label="Training Recall")
    plt.plot(range(NUM_EPOCHS), list_val_recall, label="Validation Recall")
    plt.xlabel("Epochs")
    plt.ylabel("Recall")
    plt.legend()

    plt.subplot(2, 2, 4)
    plt.plot(range(NUM_EPOCHS), list_train_average_precision, label="Training Average Precision")
    plt.plot(range(NUM_EPOCHS), list_val_average_precision, label="Validation Average Precision")
    plt.xlabel("Epochs")
    plt.ylabel("Average Precision")
    plt.legend()

    plt.tight_layout()
    save_object.savefig()
    plt.close()


def print_all_plots_to_pdf(
    save_object,
    dict_evaluation_results: dict,
    kfold_train_dataset: datasets,
    kfold_val_dataset: datasets,
):
    """
    Generates and saves various evaluation plots (e.g., class distribution, loss/accuracy curves, ROC, and precision-recall curves) to a PDF.

    Args:
        save_object (PdfPages): The object to save the plots to a PDF.
        dict_evaluation_results (dict): A dictionary containing evaluation results such as losses, accuracies, labels, and predicted probabilities.
        kfold_train_dataset (Dataset): The training dataset used for plotting class distributions.
        kfold_val_dataset (Dataset): The validation dataset used for plotting class distributions.
    """
    plot_class_distribution(
        kfold_train_dataset.__dict__["targets"], save_object, title="Training Class Distribution"
    )
    plot_class_distribution(
        kfold_val_dataset.__dict__["targets"], save_object, title="Validation Class Distribution"
    )
    plot_loss_accuracy_curve(
        dict_evaluation_results["list_train_losses"],
        dict_evaluation_results["list_val_losses"],
        dict_evaluation_results["list_train_accuracies"],
        dict_evaluation_results["list_val_accuracies"],
        save_object,
        title="Training and Validation Loss and Accuracy",
    )
    plot_roc_curve(
        dict_evaluation_results["train_all_labels"],
        dict_evaluation_results["train_all_pred_probs"],
        save_object,
        title="Training ROC Curve",
    )
    plot_roc_curve(
        dict_evaluation_results["val_all_labels"],
        dict_evaluation_results["val_all_pred_probs"],
        save_object,
        title="Validation ROC Curve",
    )
    plot_precision_recall_curve(
        dict_evaluation_results["train_all_labels"],
        dict_evaluation_results["train_all_pred_probs"],
        save_object,
        title="Training Precision-Recall Curve",
    )
    plot_precision_recall_curve(
        dict_evaluation_results["val_all_labels"],
        dict_evaluation_results["val_all_pred_probs"],
        save_object,
        title="Validation Precision-Recall Curve",
    )
    plot_train_and_val_metrics(
        save_object,
        dict_evaluation_results["list_train_f1"],
        dict_evaluation_results["list_val_f1"],
        dict_evaluation_results["list_train_precision"],
        dict_evaluation_results["list_val_precision"],
        dict_evaluation_results["list_train_recall"],
        dict_evaluation_results["list_val_recall"],
        dict_evaluation_results["list_train_average_precision"],
        dict_evaluation_results["list_val_average_precision"],
    )


def train_and_validate_fold(run, idx_fold, path_kfold_train_dir, path_kfold_val_dir):
    """
    Trains and validates the model for a single fold of cross-validation, saving results and plots.

    Args:
        run (int): The current run number.
        idx_fold (int): The current fold number in the k-fold cross-validation.
        path_kfold_train_dir (str): Path to the training dataset for the current fold.
        path_kfold_val_dir (str): Path to the validation dataset for the current fold.
    """
    run_output_dir, fold_output_dir = create_directories(run, idx_fold)
    writer = SummaryWriter(os.path.join(run_output_dir, "tensorboard"))
    with PdfPages(os.path.join(run_output_dir, f"training_report_run_{run}.pdf")) as pdf_pages:

        train_fold_loader, val_fold_loader, kfold_train_dataset, kfold_val_dataset = (
            create_dataloaders(path_kfold_train_dir, path_kfold_val_dir)
        )

        model = load_model(kfold_train_dataset)
        model.train()
        model.to(device)

        criterion = nn.CrossEntropyLoss()
        optimizer = optim.Adam(model.parameters(), lr=LR_INITIAL)
        scaler = torch.amp.GradScaler()

        list_train_losses, list_val_losses = [], []
        list_train_f1, list_val_f1 = [], []
        list_train_precision, list_val_precision = [], []
        list_train_recall, list_val_recall = [], []
        list_train_average_precision, list_val_average_precision = [], []

        highest_val_acc = 0.0
        for epoch in range(NUM_EPOCHS):
            # Training
            (
                train_running_loss,
                train_correct_preds,
                train_total_preds,
                train_all_preds,
                train_all_labels,
                train_all_pred_probs,
            ) = train_single_epoch(model, train_fold_loader, optimizer, criterion, scaler, device)
            train_loss = train_running_loss / len(train_fold_loader)
            train_acc = 100 * train_correct_preds / train_total_preds
            train_f1, train_precision, train_recall, train_average_precision = compute_metrics(
                train_all_labels, train_all_preds, train_all_pred_probs
            )

            # Validation
            (
                val_running_loss,
                val_correct_preds,
                val_total_preds,
                val_all_preds,
                val_all_labels,
                val_all_pred_probs,
            ) = validate_single_epoch(model, val_fold_loader, criterion, device)
            val_loss = val_running_loss / len(val_fold_loader)
            val_acc = 100 * val_correct_preds / val_total_preds
            val_f1, val_precision, val_recall, val_average_precision = compute_metrics(
                val_all_labels, val_all_preds, val_all_pred_probs
            )

            # Log metrics
            writer.add_scalar("Loss/train", train_loss, epoch)
            writer.add_scalar("Accuracy/train", train_acc, epoch)
            writer.add_scalar("F1/train", train_f1, epoch)
            writer.add_scalar("Precision/train", train_precision, epoch)
            writer.add_scalar("Recall/train", train_recall, epoch)
            writer.add_scalar("Average Precision/train", train_average_precision, epoch)

            writer.add_scalar("Loss/val", val_loss, epoch)
            writer.add_scalar("Accuracy/val", val_acc, epoch)
            writer.add_scalar("F1/val", val_f1, epoch)
            writer.add_scalar("Precision/val", val_precision, epoch)
            writer.add_scalar("Recall/val", val_recall, epoch)
            writer.add_scalar("Average Precision/val", val_average_precision, epoch)

            # Save best model
            save_model(model, fold_output_dir, epoch, val_acc, highest_val_acc)

        dict_evaluation_results = {
            "list_train_losses": list_train_losses,
            "list_val_losses": list_val_losses,
            "list_train_f1": list_train_f1,
            "list_val_f1": list_val_f1,
            "list_train_precision": list_train_precision,
            "list_val_precision": list_val_precision,
            "list_train_recall": list_train_recall,
            "list_val_recall": list_val_recall,
            "list_train_average_precision": list_train_average_precision,
            "list_val_average_precision": list_val_average_precision,
            "train_all_labels": train_all_labels,
            "train_all_pred_probs": train_all_pred_probs,
            "val_all_labels": val_all_labels,
            "val_all_pred_probs": val_all_pred_probs,
        }

        print_all_plots_to_pdf(
            pdf_pages, dict_evaluation_results, kfold_train_dataset, kfold_val_dataset
        )

        writer.flush()
        writer.close()


def run_training_process():
    """Main training loop for all runs."""
    for run in range(NUM_RUNS):
        print(f"\nTRAIN RUN NUMBER: {run}/{NUM_RUNS-1}")
        for idx_fold, (path_kfold_train_dir, path_kfold_val_dir) in enumerate(
            zip(list_paths_fold_train_dirs, list_paths_fold_val_dirs)
        ):
            print(f"\n----Fold {idx_fold}----------------------------")
            train_and_validate_fold(run, idx_fold, path_kfold_train_dir, path_kfold_val_dir)
        print(f"Finished run {run + 1}/{NUM_RUNS}.")


# Call the main function to start the training process
run_training_process()
