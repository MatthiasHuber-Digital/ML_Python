import matplotlib.pyplot as plt
import numpy as np
from pptx.util import Inches
from sklearn.preprocessing import LabelBinarizer
from sklearn.metrics import (
    precision_recall_curve,
    average_precision_score,
    roc_curve,
    auc,
)


def change_matplotlib_font_params():
    SMALL_SIZE = 8
    MEDIUM_SIZE = 10
    # BIGGER_SIZE = 12

    plt.rc("font", size=SMALL_SIZE)  # controls default text sizes
    plt.rc("axes", titlesize=SMALL_SIZE)  # fontsize of the axes title
    plt.rc("axes", labelsize=MEDIUM_SIZE)  # fontsize of the x and y labels
    plt.rc("xtick", labelsize=SMALL_SIZE)  # fontsize of the tick labels
    plt.rc("ytick", labelsize=SMALL_SIZE)  # fontsize of the tick labels
    plt.rc("legend", fontsize=SMALL_SIZE)  # legend fontsize
    plt.rc("figure", titlesize=MEDIUM_SIZE)  # fontsize of the figure title


def plot_roc_curve(
    y_true,
    y_pred_probs,
    save_object,
    title="ROC Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # Convert labels to one-hot format for multiclass classification
    lb = LabelBinarizer()
    y_true_bin = lb.fit_transform(y_true)

    if not isinstance(y_true_bin, np.ndarray):
        y_true_bin = np.array(y_true_bin)
    if not isinstance(y_pred_probs, np.ndarray):
        y_pred_probs = np.array(y_pred_probs)

    # Ensure y_pred_probs is a NumPy array
    y_pred_probs = np.array(y_pred_probs)

    # Check if y_pred_probs is 2D (for multiclass)
    if y_pred_probs.ndim == 1:
        # If it's 1D, expand it to 2D for binary classification (one class vs rest)
        y_pred_probs = np.expand_dims(y_pred_probs, axis=1)

    # Plot ROC curve for each class
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    n_classes = y_true_bin.shape[1]
    fpr, tpr, roc_auc = {}, {}, {}

    for i in range(n_classes):
        fpr[i], tpr[i], _ = roc_curve(y_true_bin[:, i], y_pred_probs[:, i])
        roc_auc[i] = auc(fpr[i], tpr[i])
        plt.plot(fpr[i], tpr[i], label=f"Class {lb.classes_[i]} (AUC = {roc_auc[i]:.2f})")

    plt.plot([0, 1], [0, 1], color="gray", linestyle="--")  # Diagonal line (chance level)
    plt.xlabel("False Positive Rate")
    plt.xlim([0.0, 1.0])
    plt.ylim([0.0, 1.0])
    plt.grid(color="gray", linestyle="--", linewidth=0.5)
    plt.ylabel("True Positive Rate")
    plt.title(title)
    plt.legend(loc="lower right")

    # Save the plot to the provided PDF
    save_object.savefig()
    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_precision_recall_curve(
    y_true,
    y_pred_probs,
    save_object,
    title="Precision-Recall Curve",
    # tensorboard_writer: SummaryWriter = None,
):
    # Convert labels to one-hot format for multiclass classification
    lb = LabelBinarizer()
    y_true_bin = lb.fit_transform(y_true)

    if not isinstance(y_true_bin, np.ndarray):
        y_true_bin = np.array(y_true_bin)
    if not isinstance(y_pred_probs, np.ndarray):
        y_pred_probs = np.array(y_pred_probs)

    # Ensure y_pred_probs is a NumPy array
    y_pred_probs = np.array(y_pred_probs)

    # Check if y_pred_probs is 2D (for multiclass)
    if y_pred_probs.ndim == 1:
        # If it's 1D, expand it to 2D for binary classification (one class vs rest)
        y_pred_probs = np.expand_dims(y_pred_probs, axis=1)

    # Plot Precision-Recall curve for each class
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    n_classes = y_true_bin.shape[1]

    for i in range(n_classes):
        precision, recall, _ = precision_recall_curve(y_true_bin[:, i], y_pred_probs[:, i])
        avg_precision = average_precision_score(y_true_bin[:, i], y_pred_probs[:, i])
        plt.plot(recall, precision, label=f"Class {lb.classes_[i]} (AP = {avg_precision:.2f})")

    plt.xlabel("Recall")
    plt.ylabel("Precision")
    plt.title(title)
    plt.xlim([0.0, 1.0])
    plt.ylim([0.0, 1.0])
    plt.grid(color="gray", linestyle="--", linewidth=0.5)
    plt.legend(loc="lower left")

    # Save the plot to the provided PDF
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)

    plt.close()


def plot_loss_curve(
    losses,
    save_object,
    title="Loss Curve",  # tensorboard_writer: SummaryWriter = None
):
    # fig = plt.figure(figsize=(8, 6))
    plt.figure(figsize=(8, 6))
    plt.plot(losses, color="red", lw=2)
    plt.xlabel("Epochs")
    plt.ylabel("Loss")
    plt.title(title)
    save_object.savefig()
    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_class_distribution(
    list_class_occurences,
    save_object,
    title="Class Distribution",
    # tensorboard_writer: SummaryWriter = None,
):
    # sns.barplot(x=list(class_counts.keys()), y=list(class_counts.values()))
    # fixed bin size
    bins = np.arange(0, max(list_class_occurences) + 1, 1)  # fixed bin size

    plt.xlim([min(list_class_occurences), max(list_class_occurences) + 1])
    # fig = plt.hist(list_class_occurences, bins=bins, alpha=0.5)
    plt.hist(list_class_occurences, bins=bins, alpha=0.5)
    plt.xlabel("Classes")
    plt.ylabel("Count")
    plt.title(title)
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


def plot_loss_accuracy_curve(
    num_epochs,
    list_train_losses,
    list_val_losses,
    list_train_accuracies,
    list_val_accuracies,
    save_object,
    titles: list,
    # tensorboard_writer: SummaryWriter = None,
):
    # fig = plt.figure(figsize=(10, 5))
    # plt.figure(figsize=(10, 5))
    plt.subplot(1, 2, 1)
    plt.title(titles[0])
    plt.plot(range(num_epochs), list_train_losses, label="Training Loss")
    plt.plot(range(num_epochs), list_val_losses, label="Validation Loss")
    plt.xlabel("Epochs")
    plt.ylabel("Loss")
    plt.xlim(0.0)
    plt.ylim(0.0)
    plt.grid(color="gray", linestyle="--", linewidth=0.5)
    plt.legend()

    plt.subplot(1, 2, 2)
    plt.title(titles[1])
    plt.plot(range(num_epochs), list_train_accuracies, label="Training Accuracy")
    plt.plot(range(num_epochs), list_val_accuracies, label="Validation Accuracy")
    plt.xlabel("Epochs")
    plt.ylabel("Accuracy (%)")
    plt.xlim(0.0)
    plt.ylim([0, 100])
    plt.grid(color="gray", linestyle="--", linewidth=0.5)
    plt.legend(loc="lower right")

    plt.tight_layout()
    save_object.savefig()

    # plt.show()
    # if tensorboard_writer is not None:
    #    tensorboard_writer.add_figure(title, fig, 0)
    plt.close()


# Helper function to create a PowerPoint presentation slide
def add_ppt_slide(ppt, title, content, image_path=None):
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Use the 'Title and Content' layout

    # Set slide title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Add content (text or image)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(9))  # Add image
    else:
        content_box = slide.shapes.placeholders[1]  # Text box
        content_box.text = content
